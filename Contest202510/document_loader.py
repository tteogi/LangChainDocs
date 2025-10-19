from typing import List
from langchain_community.document_loaders import PyPDFLoader, UnstructuredMarkdownLoader
from langchain.schema import Document
from pptx import Presentation
import tempfile
import os


def load_documents(uploaded_files) -> List[Document]:
    documents = []
    
    for uploaded_file in uploaded_files:
        file_extension = os.path.splitext(uploaded_file.name)[1].lower()
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as tmp_file:
            tmp_file.write(uploaded_file.getvalue())
            tmp_file_path = tmp_file.name
        
        try:
            if file_extension == '.pdf':
                loader = PyPDFLoader(tmp_file_path)
                loaded_docs = loader.load()
                for doc in loaded_docs:
                    doc.metadata["source"] = uploaded_file.name
                documents.extend(loaded_docs)
            elif file_extension in ['.md', '.markdown']:
                loader = UnstructuredMarkdownLoader(tmp_file_path)
                loaded_docs = loader.load()
                for doc in loaded_docs:
                    doc.metadata["source"] = uploaded_file.name
                documents.extend(loaded_docs)
            elif file_extension == '.pptx':
                prs = Presentation(tmp_file_path)
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = ""
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_text += shape.text + "\n"
                    
                    if slide_text.strip():
                        doc = Document(
                            page_content=slide_text.strip(),
                            metadata={
                                "source": uploaded_file.name,
                                "page": slide_num - 1,
                                "slide": slide_num
                            }
                        )
                        documents.append(doc)
        finally:
            os.unlink(tmp_file_path)
    
    return documents
